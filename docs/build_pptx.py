"""
Build ATM-Net++ presentation as a .pptx file.
Run: py -3 build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colour palette ──────────────────────────────────────────────────
BG          = RGBColor(0x0a, 0x0e, 0x1a)
SURFACE     = RGBColor(0x11, 0x18, 0x27)
ACCENT_BLUE = RGBColor(0x3b, 0x82, 0xf6)
ACCENT_PUR  = RGBColor(0x8b, 0x5c, 0xf6)
ACCENT_GRN  = RGBColor(0x10, 0xb9, 0x81)
ACCENT_AMB  = RGBColor(0xf5, 0x9e, 0x0b)
ACCENT_RED  = RGBColor(0xef, 0x44, 0x44)
WHITE       = RGBColor(0xff, 0xff, 0xff)
TEXT        = RGBColor(0xe2, 0xe8, 0xf0)
MUTED       = RGBColor(0x94, 0xa3, 0xb8)
BORDER      = RGBColor(0x1e, 0x2d, 0x45)
CARD_BG     = RGBColor(0x1a, 0x24, 0x3a)

# Slide size: 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ════════════════════════════════════════════════════════════════════
# Helper utilities
# ════════════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h,
        fill_color=CARD_BG, border_color=BORDER, border_pt=1.0,
        radius=None):
    """Add a rounded rectangle shape."""
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,   # MSO_SHAPE.RECTANGLE
        l, t, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_pt)
    return shape


def accent_bar(slide, color1=ACCENT_BLUE, color2=ACCENT_PUR):
    """Top gradient bar (simulated with two rectangles)."""
    half = SLIDE_W // 2
    r1 = slide.shapes.add_shape(1, 0, 0, half, Pt(6))
    r1.fill.solid(); r1.fill.fore_color.rgb = color1
    r1.line.fill.background()
    r2 = slide.shapes.add_shape(1, half, 0, half, Pt(6))
    r2.fill.solid(); r2.fill.fore_color.rgb = color2
    r2.line.fill.background()


def txb(slide, text, l, t, w, h,
        size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    """Add a text box."""
    tf = slide.shapes.add_textbox(l, t, w, h)
    frame = tf.text_frame
    frame.word_wrap = wrap
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def label_txt(slide, text, l, t, w=Inches(8)):
    txb(slide, text.upper(), l, t, w, Inches(0.3),
        size=9, bold=True, color=ACCENT_BLUE)


def heading(slide, text, l, t, w=Inches(12), size=32):
    txb(slide, text, l, t, w, Inches(0.7),
        size=size, bold=True, color=WHITE)


def subheading(slide, text, l, t, w=Inches(12), size=20):
    txb(slide, text, l, t, w, Inches(0.5),
        size=size, bold=True, color=TEXT)


def body_text(slide, lines, l, t, w, h, size=13, color=MUTED, bullet=True):
    """Multi-line body text with optional bullets."""
    tf = slide.shapes.add_textbox(l, t, w, h)
    frame = tf.text_frame
    frame.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = frame.paragraphs[0]
            first = False
        else:
            p = frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = ("• " if bullet else "") + line
        run.font.size = Pt(size)
        run.font.color.rgb = color


def card_with_title(slide, title, lines, l, t, w, h,
                    accent=ACCENT_BLUE, size=12):
    r = box(slide, l, t, w, h,
            fill_color=RGBColor(0x0e, 0x18, 0x2e),
            border_color=accent)
    # title bar inside card
    title_bar = slide.shapes.add_shape(1, l, t, w, Pt(28))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(
        min(accent[0]+10, 255),
        min(accent[1]+10, 255),
        min(accent[2]+10, 255))
    title_bar.line.fill.background()

    txb(slide, title,
        l + Inches(0.1), t + Pt(4),
        w - Inches(0.2), Pt(22),
        size=11, bold=True, color=WHITE)

    body_text(slide, lines,
              l + Inches(0.12), t + Pt(32),
              w - Inches(0.24), h - Pt(36),
              size=size, color=TEXT)


def slide_number(slide, n, total=14):
    txb(slide, f"{n} / {total}",
        SLIDE_W - Inches(1.1), SLIDE_H - Inches(0.35),
        Inches(1.0), Inches(0.3),
        size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def chip(slide, text, l, t, color=ACCENT_BLUE):
    w = Inches(len(text) * 0.095 + 0.25)
    h = Inches(0.28)
    r = slide.shapes.add_shape(1, l, t, w, h)
    r.fill.solid()
    # dark tinted background derived from color tuple
    cr, cg, cb = color[0], color[1], color[2]
    r.fill.fore_color.rgb = RGBColor(min(cr, 40), min(cg, 40), min(cb, 60))
    r.line.color.rgb = color
    r.line.width = Pt(0.75)
    txb(slide, text, l + Inches(0.06), t + Pt(2),
        w, h, size=9, bold=True, color=color, align=PP_ALIGN.LEFT)
    return w


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 1)

label_txt(s, "Research Project · v1.0.0", Inches(0.6), Inches(1.1))
heading(s, "ATM-Net++", Inches(0.6), Inches(1.45), size=54)
txb(s, "Anatomy-Aware Multimodal Lumbar Spine MRI",
    Inches(0.6), Inches(2.55), Inches(9), Inches(0.5),
    size=22, color=MUTED)
txb(s, "Diagnostic & Segmentation System",
    Inches(0.6), Inches(3.0), Inches(9), Inches(0.5),
    size=22, color=MUTED)

# horizontal rule
r = s.shapes.add_shape(1, Inches(0.6), Inches(3.65), Inches(5), Pt(2))
r.fill.solid(); r.fill.fore_color.rgb = BORDER; r.line.fill.background()

chips = [
    ("MRI Segmentation",      ACCENT_BLUE),
    ("Disease Classification", ACCENT_PUR),
    ("Explainability",         ACCENT_GRN),
    ("Clinical PDF Reports",   ACCENT_AMB),
    ("SPIDER Dataset",         ACCENT_RED),
]
cx = Inches(0.6)
for text, col in chips:
    w = chip(s, text, cx, Inches(4.0), col)
    cx += w + Inches(0.15)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 2)
label_txt(s, "Motivation", Inches(0.5), Inches(0.2))
heading(s, "The Clinical Problem", Inches(0.5), Inches(0.5))

card_with_title(s, "🏥  Manual Interpretation is Slow",
    ["Radiologists trace vertebrae, IVDs, and spinal canal across dozens",
     "of MRI slices — time-consuming, subjective, high inter-reader variability."],
    Inches(0.4), Inches(1.3), Inches(6.0), Inches(1.55), ACCENT_BLUE)

card_with_title(s, "📊  Multi-modal Data Underused",
    ["Clinical reports, demographics, and MRI images are rarely fused",
     "systematically, leaving valuable diagnostic signal untapped."],
    Inches(0.4), Inches(3.0), Inches(6.0), Inches(1.55), ACCENT_AMB)

card_with_title(s, "🎯  What We Need",
    ["Pixel-level segmentation of vertebrae & discs",
     "Automatic disease detection & severity grading",
     "Explainable, trustworthy AI predictions",
     "Fusion of MRI + text reports + demographics",
     "Clinician-ready PDF output"],
    Inches(6.7), Inches(1.3), Inches(6.0), Inches(1.9), ACCENT_GRN)

card_with_title(s, "💡  ATM-Net++ Solution",
    ["Production-grade clinical AI platform combining MRI, radiology",
     "text and patient demographics to deliver segmentation maps,",
     "disease labels, severity scores, and structured reports —",
     "all with Grad-CAM explainability."],
    Inches(6.7), Inches(3.35), Inches(6.0), Inches(1.8), ACCENT_PUR)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Key Numbers
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 3)
label_txt(s, "At a Glance", Inches(0.5), Inches(0.2))
heading(s, "Key Numbers", Inches(0.5), Inches(0.5))

stats_top = [
    (">0.90",  "Target Dice Score",     ACCENT_BLUE),
    ("20",     "Segmentation Classes",  ACCENT_GRN),
    ("7",      "Disease Categories",    ACCENT_PUR),
    ("447",    "SPIDER MRI Volumes",    ACCENT_AMB),
]
stats_bot = [
    ("3",      "Input Modalities",      ACCENT_BLUE),
    ("5",      "Docker Micro-services", ACCENT_GRN),
    ("200",    "Training Epochs",       ACCENT_PUR),
    ("5-fold", "Cross-Validation",      ACCENT_AMB),
]

CW = Inches(2.9); CH = Inches(1.8); GAP = Inches(0.3)
start_x = Inches(0.5)
for i, (num, desc, col) in enumerate(stats_top):
    lx = start_x + i * (CW + GAP)
    r = box(s, lx, Inches(1.4), CW, CH,
            fill_color=RGBColor(0x0e, 0x18, 0x2e), border_color=col)
    txb(s, num,  lx, Inches(1.65), CW, Inches(0.65),
        size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, desc, lx, Inches(2.4),  CW, Inches(0.4),
        size=12, color=MUTED, align=PP_ALIGN.CENTER)

for i, (num, desc, col) in enumerate(stats_bot):
    lx = start_x + i * (CW + GAP)
    r = box(s, lx, Inches(3.45), CW, Inches(1.55),
            fill_color=RGBColor(0x0e, 0x18, 0x2e), border_color=col)
    txb(s, num,  lx, Inches(3.62), CW, Inches(0.55),
        size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, desc, lx, Inches(4.25), CW, Inches(0.35),
        size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dataset
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 4)
label_txt(s, "Data", Inches(0.5), Inches(0.2))
heading(s, "SPIDER Dataset", Inches(0.5), Inches(0.5))

card_with_title(s, "📦  Dataset Overview",
    ["447 MRI volumes (T1 & T2 sagittal)",
     "Paired images + masks in .mha format",
     "overview.csv — acquisition metadata + sex/subset",
     "radiological_gradings.csv — per-IVD pathology grades"],
    Inches(0.4), Inches(1.3), Inches(5.8), Inches(1.75), ACCENT_BLUE)

card_with_title(s, "🔀  Dataset Split",
    ["70% Train  ·  15% Validation  ·  15% Test",
     "5-fold cross-validation strategy",
     "Both T1 and T2 modalities used",
     "Image resolution: 512 × 512 pixels"],
    Inches(0.4), Inches(3.2), Inches(5.8), Inches(1.75), ACCENT_AMB)

# Label mapping table (right column)
box(s, Inches(6.5), Inches(1.3), Inches(6.5), Inches(3.7),
    fill_color=RGBColor(0x0e, 0x18, 0x2e), border_color=ACCENT_PUR)
txb(s, "Label Mapping", Inches(6.6), Inches(1.35), Inches(6.3), Inches(0.35),
    size=11, bold=True, color=WHITE)

rows = [
    ("SPIDER ID", "Structure",           "ATM-Net++ ID"),
    ("—",         "Background",           "0"),
    ("20–24",     "L1 – L5 vertebrae",    "4 – 8"),
    ("25",        "S1",                   "9"),
    ("119–122",   "L1/L2 – L4/L5 discs", "13 – 16"),
    ("123",       "L5/S1 disc",           "17"),
    ("201",       "Spinal canal",         "18"),
    ("—",         "Spinal cord",          "19"),
]
col_x = [Inches(6.55), Inches(8.2), Inches(11.2)]
row_h = Inches(0.36)
for ri, row in enumerate(rows):
    ty = Inches(1.72) + ri * row_h
    bg_col = RGBColor(0x1a, 0x26, 0x40) if ri == 0 else RGBColor(0x0e, 0x18, 0x2e)
    for ci, cell in enumerate(row):
        txb(s, cell, col_x[ci], ty, Inches(2.5), row_h,
            size=11,
            bold=(ri == 0),
            color=WHITE if ri == 0 else (ACCENT_BLUE if ci == 2 else MUTED))


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — System Architecture
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 5)
label_txt(s, "Design", Inches(0.5), Inches(0.2))
heading(s, "System Architecture", Inches(0.5), Inches(0.5))

# Column positions
cols_x = [Inches(0.3), Inches(3.1), Inches(6.3), Inches(9.5), Inches(11.8)]
arrows_x = [Inches(2.6), Inches(5.75), Inches(8.95), Inches(11.35)]
item_h = Inches(1.1)
item_w = Inches(2.6)

sections = [
    ("Inputs",    ACCENT_BLUE,
     ["MRI Image\n(T1/T2, 512×512)", "Radiology Text\n(report)", "Demographics\n(age, sex, BMI)"]),
    ("Encoders",  ACCENT_PUR,
     ["Swin UNETR\n768-dim", "Bio-ClinicalBERT\nCLS 768-dim", "Demographic MLP\n256-dim"]),
    ("Fusion",    ACCENT_GRN,
     ["Multimodal Fusion\nATPG · HASF · CCAE\n4 layers · 8 heads"]),
    ("Outputs",   ACCENT_RED,
     ["Segmentation\n20-class mask", "Disease Class\n7 categories", "Severity + Grade\nMild/Mod/Severe", "PDF Report\n+ Grad-CAM"]),
]

label_y = Inches(1.1)
for si, (sec_name, col, items) in enumerate(sections):
    lx = cols_x[si]
    txb(s, sec_name, lx, label_y, item_w, Inches(0.3),
        size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    for ii, item in enumerate(items):
        ty = Inches(1.5) + ii * (item_h + Inches(0.12))
        r = box(s, lx, ty, item_w, item_h,
                fill_color=RGBColor(0x0e, 0x18, 0x2e), border_color=col)
        txb(s, item, lx, ty + Inches(0.2), item_w, item_h,
            size=11, color=col, align=PP_ALIGN.CENTER, bold=False)

    # arrow after each section except last
    if si < len(sections) - 1:
        ax = cols_x[si] + item_w + Inches(0.05)
        txb(s, "→", ax, Inches(2.5), Inches(0.35), Inches(0.5),
            size=22, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Novel Contributions
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 6)
label_txt(s, "Novel Contributions", Inches(0.5), Inches(0.2))
heading(s, "What Makes ATM-Net++ New", Inches(0.5), Inches(0.5))

contribs = [
    ("① ATPG — Anatomy-Text Prompt Generation",
     ["Generates anatomy-aware text prompts guided by image features.",
      "Bridges visual and language spaces by conditioning BERT token",
      "attention on spatial encoder activations.",
      "16 learnable prompts  ·  Cross-modal conditioning"],
     ACCENT_BLUE, Inches(0.4), Inches(1.3)),
    ("② HASF — Hierarchical Anatomy-Aware Semantic Fusion",
     ["Fuses image, text and demographic embeddings through multi-head",
      "cross-attention at multiple hierarchical scales.",
      "Preserves both global context and local anatomy detail.",
      "8 attention heads  ·  4 transformer layers"],
     ACCENT_PUR, Inches(6.8), Inches(1.3)),
    ("③ CCAE — Cross-modal Context-Aware Enhancement",
     ["FiLM-style modulation: injects fused multimodal context into decoder",
      "feature maps via learned scale and shift parameters.",
      "Refines segmentation using semantic knowledge from text & demographics.",
      "FiLM conditioning  ·  Decoder injection"],
     ACCENT_GRN, Inches(0.4), Inches(4.0)),
    ("④ Joint Contrastive Image-Text Alignment",
     ["NT-Xent contrastive loss aligns image and text feature projections",
      "in a shared 256-dim embedding space.",
      "Improves cross-modal grounding without extra labeled pairs.",
      "NT-Xent loss  ·  256-dim shared space"],
     ACCENT_AMB, Inches(6.8), Inches(4.0)),
]

for title, lines, col, lx, ty in contribs:
    card_with_title(s, title, lines, lx, ty, Inches(6.1), Inches(2.35), col, size=12)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Model Components
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 7)
label_txt(s, "Implementation", Inches(0.5), Inches(0.2))
heading(s, "Model Components", Inches(0.5), Inches(0.5))

components = [
    ("🖼  Image Encoder",
     ["Swin UNETR backbone (MONAI)",
      "Residual blocks + attention gates",
      "Deep supervision (3 levels)",
      "Feature size 48 → 768-dim global",
      "Lightweight U-Net fallback"],
     ACCENT_BLUE),
    ("📝  Text Encoder",
     ["Bio-ClinicalBERT (Alsentzer 2019)",
      "First 6 layers frozen",
      "Fine-tuned classification head",
      "CLS embedding: 768-dim",
      "Max sequence: 512 tokens"],
     ACCENT_PUR),
    ("👤  Demographic Encoder",
     ["8-dim input vector",
      "3-layer MLP: 64 → 128 → 256",
      "Dropout 0.2",
      "Encodes age, sex, BMI, height,",
      "weight, symptoms"],
     ACCENT_GRN),
    ("🧩  Multi-Task Heads",
     ["Disease classification (7-class)",
      "Severity: 3-class + regression",
      "Level localization (8 IVDs)",
      "IVD pathology + Pfirrmann grade",
      "All on fused 512-dim features"],
     ACCENT_AMB),
    ("📄  Report Generator",
     ["NeuralReportHead",
      "Template + neural hybrid",
      "Findings, impression, recommendation",
      "ReportLab PDF rendering",
      "Radiologist-style structure"],
     ACCENT_RED),
    ("🔍  Explainability",
     ["Grad-CAM (Selvaraju 2017)",
      "Attention Rollout (Abnar 2020)",
      "Segmentation overlay (α=0.4)",
      "Disease localization heatmaps",
      "Base64 encoding for API"],
     RGBColor(0x06, 0xb6, 0xd4)),
]

CW2 = Inches(4.1); CH2 = Inches(2.4); GAPX = Inches(0.3); GAPY = Inches(0.2)
for i, (title, lines, col) in enumerate(components):
    col_i = i % 3
    row_i = i // 3
    lx = Inches(0.35) + col_i * (CW2 + GAPX)
    ty = Inches(1.3)  + row_i * (CH2 + GAPY)
    card_with_title(s, title, lines, lx, ty, CW2, CH2, col, size=11)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Training Details
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 8)
label_txt(s, "Training", Inches(0.5), Inches(0.2))
heading(s, "Training Configuration & Losses", Inches(0.5), Inches(0.5))

card_with_title(s, "⚙️  Optimizer & Schedule",
    ["AdamW  ·  lr=1e-4  ·  weight_decay=1e-5  ·  β=(0.9, 0.999)",
     "Cosine LR with 10-epoch linear warmup, min LR=1e-6",
     "200 epochs  ·  batch size 4  ·  grad accumulation ×4",
     "Mixed precision FP16 (torch.cuda.amp)",
     "Early stopping: patience=30 on val Dice"],
    Inches(0.4), Inches(1.3), Inches(6.1), Inches(2.0), ACCENT_BLUE)

card_with_title(s, "🔁  Data Augmentation",
    ["Random rotation ±15°",
     "Horizontal/vertical flip (p=0.5)",
     "Elastic deformation (α=100, σ=10)",
     "Intensity shift ±0.1, scale ±0.1",
     "Gaussian noise σ=0.01",
     "Random crop 480×480, zoom 0.9–1.1×"],
    Inches(0.4), Inches(3.45), Inches(6.1), Inches(2.3), ACCENT_AMB)

# Loss table
box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(4.4),
    fill_color=RGBColor(0x0e, 0x18, 0x2e), border_color=ACCENT_GRN)
txb(s, "📉  Loss Functions", Inches(6.9), Inches(1.35), Inches(6.0), Inches(0.35),
    size=11, bold=True, color=WHITE)

loss_rows = [
    ("Loss Component",                "Weight"),
    ("Segmentation Dice",             "1.0"),
    ("Segmentation Focal",            "0.5"),
    ("Segmentation Boundary",         "0.2"),
    ("Disease Classification (Focal CE)", "0.3"),
    ("Severity Estimation",           "0.2"),
    ("Level Localization",            "0.2"),
    ("Report Generation",             "0.1"),
    ("Contrastive NT-Xent",           "0.1"),
]
for ri, (comp, wt) in enumerate(loss_rows):
    ty = Inches(1.72) + ri * Inches(0.4)
    txb(s, comp, Inches(6.95), ty, Inches(4.5), Inches(0.38),
        size=11, bold=(ri==0), color=WHITE if ri==0 else MUTED)
    txb(s, wt,   Inches(11.55), ty, Inches(1.2), Inches(0.38),
        size=11, bold=(ri==0), color=ACCENT_GRN if ri>0 else WHITE,
        align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — Evaluation
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 9)
label_txt(s, "Evaluation", Inches(0.5), Inches(0.2))
heading(s, "Metrics & Evaluation Protocol", Inches(0.5), Inches(0.5))

card_with_title(s, "📐  Segmentation Metrics",
    ["Dice Score — per-class & mean (target > 0.90)",
     "Jaccard / IoU — intersection over union",
     "HD95 — 95th percentile Hausdorff Distance",
     "ASD — Average Surface Distance",
     "Precision, Recall, F1"],
    Inches(0.4), Inches(1.3), Inches(6.1), Inches(2.0), ACCENT_BLUE)

card_with_title(s, "🏥  Classification Metrics",
    ["Disease classification accuracy",
     "Macro F1 score across 7 classes",
     "Per-class AUC",
     "Severity estimation MAE + accuracy"],
    Inches(0.4), Inches(3.45), Inches(6.1), Inches(1.75), ACCENT_PUR)

card_with_title(s, "🔬  Inference (TTA)",
    ["Test-Time Augmentation — 3 flip combinations",
     "Threshold: 0.5 on softmax probabilities",
     "Batch size 1 for production inference"],
    Inches(6.8), Inches(1.3), Inches(6.2), Inches(1.6), ACCENT_GRN)

card_with_title(s, "💾  Output Artifacts",
    ["Segmentation mask (.npy / .mha)",
     "Segmentation overlay PNG",
     "Grad-CAM heatmap PNG",
     "Structured JSON results",
     "Clinical PDF report"],
    Inches(6.8), Inches(3.05), Inches(6.2), Inches(1.9), ACCENT_AMB)

# Disease class chips
txb(s, "Disease Classes:", Inches(0.4), Inches(5.3), Inches(3), Inches(0.3),
    size=10, bold=True, color=MUTED)
diseases = [("Normal",ACCENT_GRN),("Disc Herniation",ACCENT_AMB),
            ("Disc Bulge",ACCENT_AMB),("Spinal Stenosis",ACCENT_RED),
            ("DDD",ACCENT_RED),("Spondylolisthesis",ACCENT_RED),
            ("Compression Fracture",ACCENT_RED)]
cx = Inches(0.4)
for name, col in diseases:
    w = chip(s, name, cx, Inches(5.65), col)
    cx += w + Inches(0.12)
    if cx > Inches(12.5):
        cx = Inches(0.4)


# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — Backend API
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 10)
label_txt(s, "Backend", Inches(0.5), Inches(0.2))
heading(s, "FastAPI Backend & Database", Inches(0.5), Inches(0.5))

# API table
box(s, Inches(0.4), Inches(1.3), Inches(6.3), Inches(4.5),
    fill_color=RGBColor(0x0e, 0x18, 0x2e), border_color=ACCENT_BLUE)
txb(s, "🔌  REST API Endpoints", Inches(0.5), Inches(1.35), Inches(6.1), Inches(0.35),
    size=11, bold=True, color=WHITE)

endpoints = [
    ("Method", "Route",                       "Purpose"),
    ("POST",   "/auth/register",              "Register user"),
    ("POST",   "/auth/login",                 "Get JWT token"),
    ("POST",   "/predict/upload-mri",         "Upload & predict"),
    ("POST",   "/predict/segment",            "Segmentation only"),
    ("GET",    "/patients",                   "List patients"),
    ("GET",    "/reports/download/{id}/pdf",  "Download PDF"),
    ("GET",    "/analytics/summary",          "Dashboard stats"),
    ("GET",    "/health",                     "Health check"),
]
for ri, (meth, route, desc) in enumerate(endpoints):
    ty = Inches(1.75) + ri * Inches(0.37)
    col_m = ACCENT_GRN if meth=="GET" else ACCENT_AMB if meth=="POST" else MUTED
    txb(s, meth,  Inches(0.55), ty, Inches(0.7), Inches(0.35),
        size=10, bold=(ri==0), color=WHITE if ri==0 else col_m)
    txb(s, route, Inches(1.35), ty, Inches(3.1), Inches(0.35),
        size=10, bold=(ri==0), color=WHITE if ri==0 else ACCENT_BLUE)
    txb(s, desc,  Inches(4.55), ty, Inches(2.0), Inches(0.35),
        size=10, bold=(ri==0), color=WHITE if ri==0 else MUTED)

card_with_title(s, "🗄  Database Schema (PostgreSQL)",
    ["Users — roles: admin, radiologist, clinician, researcher, viewer",
     "Patients — demographics, medical history (JSON)",
     "Studies — MRI metadata, modality, status lifecycle",
     "Predictions — seg paths, Dice, disease, severity, levels, Grad-CAM",
     "Reports — findings, impression, PDF path, review workflow",
     "AuditLogs — HIPAA-style action tracking"],
    Inches(6.9), Inches(1.3), Inches(6.1), Inches(2.4), ACCENT_PUR)

card_with_title(s, "🔐  Security",
    ["JWT authentication (HS256, 60-min expiry)",
     "bcrypt password hashing (passlib)",
     "Role-based access control (5 roles)",
     "Async SQLAlchemy + asyncpg connection pool"],
    Inches(6.9), Inches(3.85), Inches(6.1), Inches(1.85), ACCENT_GRN)


# ════════════════════════════════════════════════════════════════════
# SLIDE 11 — Frontend
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 11)
label_txt(s, "Frontend", Inches(0.5), Inches(0.2))
heading(s, "Next.js 14 Clinical Dashboard", Inches(0.5), Inches(0.5))

# Routes table
box(s, Inches(0.4), Inches(1.3), Inches(6.1), Inches(2.4),
    fill_color=RGBColor(0x0e, 0x18, 0x2e), border_color=ACCENT_BLUE)
txb(s, "📱  Pages & Routes", Inches(0.5), Inches(1.35), Inches(5.9), Inches(0.35),
    size=11, bold=True, color=WHITE)
routes = [
    ("Route",         "Description"),
    ("/auth/login",   "JWT login page"),
    ("/dashboard",    "Analytics overview & charts"),
    ("/upload",       "Drag-drop MRI upload, inline results"),
    ("/patients",     "Patient CRUD management"),
]
for ri, (r, d) in enumerate(routes):
    ty = Inches(1.75) + ri * Inches(0.42)
    txb(s, r, Inches(0.55), ty, Inches(2.5), Inches(0.38),
        size=11, bold=(ri==0), color=WHITE if ri==0 else ACCENT_BLUE)
    txb(s, d, Inches(3.15), ty, Inches(3.2), Inches(0.38),
        size=11, bold=(ri==0), color=WHITE if ri==0 else MUTED)

card_with_title(s, "🧰  Tech Stack",
    ["Next.js 14  ·  TypeScript  ·  App Router",
     "Tailwind CSS — utility-first styling",
     "Zustand — lightweight auth state management",
     "Full TypeScript API client (src/lib/api.ts)"],
    Inches(0.4), Inches(3.85), Inches(6.1), Inches(1.85), ACCENT_AMB)

card_with_title(s, "🔧  Frontend Features",
    ["Drag-and-drop MRI upload (MHA/NIfTI/DICOM/PNG)",
     "Real-time prediction results display",
     "Segmentation overlay & Grad-CAM viewer",
     "Disease classification confidence bars",
     "PDF report download",
     "Analytics charts & summary dashboard"],
    Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.4), ACCENT_GRN)

card_with_title(s, "🌐  API Client",
    ["Type-safe api.ts covers all backend endpoints",
     "Zustand authStore manages JWT token lifecycle",
     "Axios-based with automatic auth headers",
     "Error handling & loading states built-in"],
    Inches(6.8), Inches(3.85), Inches(6.2), Inches(1.85), ACCENT_PUR)


# ════════════════════════════════════════════════════════════════════
# SLIDE 12 — Deployment
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 12)
label_txt(s, "Deployment", Inches(0.5), Inches(0.2))
heading(s, "Docker Micro-services Stack", Inches(0.5), Inches(0.5))

services = [
    ("🌐", "Nginx",      "Reverse proxy\nPort 80 · SSL termination",     ACCENT_BLUE),
    ("⚛️", "Frontend",   "Next.js · Port 3000\nProduction build",         ACCENT_PUR),
    ("⚡", "Backend",    "FastAPI / Uvicorn\nPort 8000 · 4 workers",      ACCENT_GRN),
    ("🐘", "PostgreSQL", "Port 5432\nPool size 10 · max overflow 20",      ACCENT_AMB),
    ("🔴", "Redis",      "Port 6379\nCache + Celery task queue",           ACCENT_RED),
    ("📦", "ONNX Export","Optimized inference\nscripts/export_onnx.py",   RGBColor(0x06,0xb6,0xd4)),
]

SW = Inches(2.0); SH = Inches(2.1); SGAP = Inches(0.2)
sx0 = Inches(0.35)
for i, (icon, name, desc, col) in enumerate(services):
    lx = sx0 + i * (SW + SGAP)
    r = box(s, lx, Inches(1.3), SW, SH,
            fill_color=RGBColor(0x0e, 0x18, 0x2e), border_color=col)
    txb(s, icon, lx, Inches(1.38), SW, Inches(0.5),
        size=22, align=PP_ALIGN.CENTER, color=col)
    txb(s, name, lx, Inches(1.88), SW, Inches(0.35),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, desc, lx, Inches(2.28), SW, Inches(0.9),
        size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Quick commands
box(s, Inches(0.35), Inches(3.65), Inches(12.65), Inches(2.0),
    fill_color=RGBColor(0x06, 0x0d, 0x18), border_color=BORDER)
txb(s, "docker-compose up --build -d",
    Inches(0.55), Inches(3.75), Inches(12.3), Inches(0.4),
    size=14, bold=True, color=ACCENT_GRN)
txb(s, "Frontend: http://localhost:3000     API Docs: http://localhost:8000/docs     Health: http://localhost:8000/health",
    Inches(0.55), Inches(4.2), Inches(12.3), Inches(0.35),
    size=11, color=MUTED)
txb(s, "Backend only:  docker-compose up backend db redis -d     •     Logs:  docker-compose logs -f backend",
    Inches(0.55), Inches(4.6), Inches(12.3), Inches(0.35),
    size=11, color=MUTED)


# ════════════════════════════════════════════════════════════════════
# SLIDE 13 — Research Foundation
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 13)
label_txt(s, "Research Foundation", Inches(0.5), Inches(0.2))
heading(s, "Based On & Extends", Inches(0.5), Inches(0.5))

# Prior art table
box(s, Inches(0.4), Inches(1.3), Inches(6.3), Inches(3.8),
    fill_color=RGBColor(0x0e, 0x18, 0x2e), border_color=ACCENT_BLUE)
txb(s, "📚  Prior Art", Inches(0.5), Inches(1.35), Inches(6.1), Inches(0.35),
    size=11, bold=True, color=WHITE)
refs = [
    ("Work",                        "Used For"),
    ("Swin UNETR (Liu 2021)",       "Image backbone"),
    ("Bio-ClinicalBERT (2019)",     "Text encoder"),
    ("Attention U-Net (Oktay 2018)","Attention gates"),
    ("Boundary Loss (Kervadec 2019)","Seg. loss term"),
    ("Focal Loss (Lin 2017)",       "Class imbalance"),
    ("GradCAM (Selvaraju 2017)",    "Explainability"),
    ("ATM-Net (base model)",        "Anatomy-text guidance"),
]
for ri, (w, u) in enumerate(refs):
    ty = Inches(1.75) + ri * Inches(0.38)
    txb(s, w, Inches(0.55), ty, Inches(3.9), Inches(0.36),
        size=10, bold=(ri==0), color=WHITE if ri==0 else ACCENT_BLUE)
    txb(s, u, Inches(4.55), ty, Inches(2.0), Inches(0.36),
        size=10, bold=(ri==0), color=WHITE if ri==0 else MUTED)

card_with_title(s, "🔑  Key Python Dependencies",
    ["PyTorch 2.3  ·  MONAI 1.3  ·  Transformers 4.40",
     "FastAPI 0.111  ·  SQLAlchemy 2.0  ·  asyncpg",
     "SimpleITK 2.3  ·  OpenCV 4.9  ·  Albumentations 1.4",
     "ReportLab 4.1  ·  WandB 0.17  ·  TorchMetrics 1.4",
     "Celery 5.4  ·  Redis  ·  Hydra/OmegaConf"],
    Inches(6.9), Inches(1.3), Inches(6.1), Inches(2.3), ACCENT_GRN)

card_with_title(s, "🖥  Hardware Requirements",
    ["Python 3.10+  ·  Node.js 20+",
     "8 GB+ RAM (16 GB+ recommended for training)",
     "GPU optional for inference, required for fast training",
     "Docker + Docker Compose for full stack deployment"],
    Inches(6.9), Inches(3.75), Inches(6.1), Inches(1.8), ACCENT_AMB)

card_with_title(s, "⚠️  License",
    ["Research use only.",
     "NOT for clinical deployment without additional",
     "validation and regulatory approval."],
    Inches(0.4), Inches(5.25), Inches(6.3), Inches(1.5), ACCENT_RED, size=12)


# ════════════════════════════════════════════════════════════════════
# SLIDE 14 — Quick Start
# ════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s); slide_number(s, 14)
label_txt(s, "Getting Started", Inches(0.5), Inches(0.2))
heading(s, "Quick Start", Inches(0.5), Inches(0.5))

steps = [
    ("① Install dependencies",
     'pip install -r requirements.txt\ncd frontend && npm install && cd ..',
     Inches(0.4), Inches(1.3), ACCENT_BLUE),
    ("② Configure & set up data",
     'copy .env.example .env\n# Edit SECRET_KEY and DB settings\n\npython scripts/setup_data.py \\\n  --source "C:\\project\\Spine Segmentation\\10159290"',
     Inches(0.4), Inches(3.1), ACCENT_AMB),
    ("③ Launch full stack",
     'docker-compose up --build\n\n# Frontend  → http://localhost:3000\n# API docs  → http://localhost:8000/docs',
     Inches(6.8), Inches(1.3), ACCENT_GRN),
    ("④ Train & run inference",
     'python training/train.py \\\n  --config configs/base_config.yaml\n\npython scripts/run_inference.py \\\n  --image data/.../100_t2.mha \\\n  --checkpoint checkpoints/best.pth \\\n  --save-report outputs/report.pdf',
     Inches(6.8), Inches(3.1), ACCENT_PUR),
]

for title, code, lx, ty, col in steps:
    box(s, lx, ty, Inches(6.1), Inches(1.6),
        fill_color=RGBColor(0x06, 0x0d, 0x18), border_color=col)
    txb(s, title, lx + Inches(0.12), ty + Inches(0.08),
        Inches(5.9), Inches(0.32), size=11, bold=True, color=col)
    txb(s, code, lx + Inches(0.12), ty + Inches(0.42),
        Inches(5.85), Inches(1.1), size=10, color=ACCENT_GRN)

# Final badge row
final_chips = [
    ("✔ Full Docker stack",        ACCENT_GRN),
    ("✔ Interactive API docs",     ACCENT_BLUE),
    ("✔ W&B + TensorBoard",        ACCENT_PUR),
    ("✔ ONNX export ready",        ACCENT_AMB),
    ("✔ pytest unit+integration",  ACCENT_RED),
]
cx = Inches(0.4)
for text, col in final_chips:
    w = chip(s, text, cx, Inches(6.9), col)
    cx += w + Inches(0.2)


# ════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════
out = r"c:\project\Spine Segmentation\ATM-Net++\docs\ATM_Net_PlusPlus_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
